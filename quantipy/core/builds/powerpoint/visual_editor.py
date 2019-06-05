# encoding: utf-8

'''
@author: Majeed.sahebzadha
'''


import numpy as np
import pandas as pd
import re
import time
import os
from collections import OrderedDict
import json
import pickle
from pptx import Presentation
from pptx.chart.data import ChartData
from .add_shapes import *
from .transformations import *
from os.path import ( 
    basename, 
    dirname,
    split
    )
from pptx.enum.chart import (
    XL_CHART_TYPE, 
    XL_LABEL_POSITION, 
    XL_LEGEND_POSITION, 
    XL_TICK_MARK, 
    XL_TICK_LABEL_POSITION
    )
from pptx.util import (
    Emu,
    Pt,
    Cm,
    Inches
    )
from pptx.enum.dml import (
    MSO_THEME_COLOR, 
    MSO_COLOR_TYPE,
    MSO_FILL
    )
from pptx.enum.text import (
    PP_ALIGN,
    MSO_AUTO_SIZE, 
    MSO_ANCHOR
    )


pd.set_option('display.expand_frame_repr', False)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def file_exists(file_name):

    ''' check if file exists '''
    return os.path.isfile(file_name) 

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def pickle_file(obj_to_pickle, file_name):

    if ".pickle" not in file_name:
        file_name = "%s.pickle" % file_name

    with open(file_name, 'wb') as handle:
        pickle.dump(obj_to_pickle, handle)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def load_pickled_file(file_name):

    if ".pickle" not in file_name:
        file_name = "%s.pickle" % file_name

    with open(file_name, 'rb') as handle:
        picked_obj = pickle.load(handle)

    return picked_obj

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def save_json(obj, json_path):
    ''' Saves obj as a json text file at json_path
    '''

    from decoder_ring import as_ascii

    print("Saving json: {f}".format(f=json_path))
    obj = as_ascii(obj, control=False, extended=True, encoding='UTF-8')
    with open(json_path, 'w') as f:
        json.dump(obj, f)       

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def rename_duplicate_shape_names(pres_path, overwrite=True):
    ''' 
    Ensures all shapes have a unique name. 
    Only renames duplicates. 
    Compares shape names one slide at a time. 
    '''

    file_name = basename(pres_path).split('.')[0]
    file_path = dirname(pres_path)

    prs = Presentation(pres_path)

    for slide in prs.slides:
        shape_names = []
        for shape in slide.shapes:
            shape_names.append(shape.name)
        renamed_shapes = [x + "_" + str(i) if shape_names.count(x)>1 else x for i, x in enumerate(shape_names)]
        for s_idx, shape in enumerate(slide.shapes):
            shape.name = renamed_shapes[s_idx]

    if overwrite:
        prs.save('{pres_path}\\{pres_name}.pptx'.format(
            pres_path=file_path, 
            pres_name=file_name))
    else:
        prs.save('{pres_path}\\{pres_name}_edited.pptx'.format(
            pres_path=file_path, 
            pres_name=file_name))

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def read_pptx(pres_path, slide_num=[], save_as_json=True):
    ''' 
    Iterates through an existing PPTX and prints info re slide and shapes.

    param: pres_path - full path of target file
    param: slide_num - list
    parem: save_as_json - boolean

    example useage:
    *read_pptx(pres)
    *read_pptx(pres, [20,25,15], False)
    '''

    if not isinstance(slide_num, list):
        slide_num = [slide_num]

    prs = Presentation(pres_path)
    
    file_name = os.path.basename(pres_path).split('.')[0]

    pptx_tree = OrderedDict()
    pptx_tree[file_name] = OrderedDict()
    pptx_tree[file_name]['slides'] = OrderedDict()

    print('Analysing PPTX content...\n')

    for i, sld in enumerate(prs.slides, start=1):
        if slide_num:
            if i in slide_num:

                slide_number = str(i)
                pptx_tree[file_name]['slides'][slide_number] = OrderedDict()
                
                print('{indent:>5}slide layout name : {sld_layout_name}\n'.
                    format(
                        indent='', 
                        sld_layout_name=sld.slide_layout.name))

                pptx_tree[file_name]['slides'][slide_number]['slide layout'] = OrderedDict()
                slide_layout_name = str(sld.slide_layout.name)
                pptx_tree[file_name]['slides'][slide_number]['slide layout']['name'] = slide_layout_name
                
                pptx_tree[file_name]['slides'][slide_number]['shapes'] = OrderedDict()

                for x, shp in enumerate(sld.shapes):
                    print('{indent:>10}shape index - {x}'.
                        format(
                            indent='',
                            x=x))

                    shape_number = str(x)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number] = OrderedDict()

                    print('{indent:>15}shape name - {shape_name}'.
                        format(
                            indent='', 
                            shape_name=shp.name))

                    shape_name = str(shp.name)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape name'] = shape_name

                    print('{indent:>15}shape type - {shape_type}'.
                        format(
                            indent='', 
                            shape_type=shp.shape_type))

                    shape_type = str(shp.shape_type)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape type'] = shape_type

                    if str(shp.shape_type) == 'PLACEHOLDER (14)':
                        print('{indent1:>15}placeholder idx - {placeholder_idx},\n'
                            '{indent2:>15}placeholder type - {placeholder_type}'.
                                format(
                                    indent1='', 
                                    indent2='', 
                                    placeholder_idx=shp.placeholder_format.idx, 
                                    placeholder_type=shp.placeholder_format.type))

                        pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['placeholder idx'] = str(shp.placeholder_format.idx)
                        pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['placeholder type'] = str(shp.placeholder_format.type)


                    print('{indent:>15}shape dimensions - '
                        'left: {shape_left}, top: {shape_top}, '
                        'height: {shape_height}, width: {shape_width}\n'.
                          format(
                            indent='', 
                            shape_left=shp.left, 
                            shape_height=shp.height, 
                            shape_top=shp.top, 
                            shape_width=shp.width))

                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions'] = OrderedDict()
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['left'] = str(shp.left)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['top'] = str(shp.top)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['width'] = str(shp.width)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['height'] = str(shp.height)

        else:
            print('='*110)
            print('{indent:>0}Slide {i} details:\n'.
                format(
                    indent='', 
                    i=i)) 

            slide_number = str(i)
            pptx_tree[file_name]['slides'][slide_number] = OrderedDict()
            
            print('{indent:>5}slide layout name : {sld_layout_name}\n'.
                format(
                    indent='', 
                    sld_layout_name=sld.slide_layout.name))

            pptx_tree[file_name]['slides'][slide_number]['slide layout'] = OrderedDict()
            slide_layout_name = str(sld.slide_layout.name)
            pptx_tree[file_name]['slides'][slide_number]['slide layout']['name'] = slide_layout_name
            
            pptx_tree[file_name]['slides'][slide_number]['shapes'] = OrderedDict()

            for x, shp in enumerate(sld.shapes):
                print('{indent:>10}shape index - {x}'.
                    format(
                        indent='',
                        x=x))

                shape_number = str(x)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number] = OrderedDict()

                print('{indent:>15}shape name - {shape_name}'.
                    format(
                        indent='', 
                        shape_name=shp.name))

                shape_name = str(shp.name)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape name'] = shape_name

                print('{indent:>15}shape id - {shape_id}'.
                    format(
                        indent='', 
                        shape_id=shp.id))

                shape_id = str(shp.id)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape id'] = shape_id

                print('{indent:>15}shape type - {shape_type}'.
                    format(
                        indent='', 
                        shape_type=shp.shape_type))

                shape_type = str(shp.shape_type)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape type'] = shape_type

                if str(shp.shape_type) == 'PLACEHOLDER (14)':
                    print('{indent1:>15}placeholder idx - {placeholder_idx},\n'
                        '{indent2:>15}placeholder type - {placeholder_type}'.
                            format(
                                indent1='', 
                                indent2='', 
                                placeholder_idx=shp.placeholder_format.idx, 
                                placeholder_type=shp.placeholder_format.type))

                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['placeholder idx'] = str(shp.placeholder_format.idx)
                    pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['placeholder type'] = str(shp.placeholder_format.type)


                print('{indent:>15}shape dimensions - '
                    'left: {shape_left}, top: {shape_top}, '
                    'height: {shape_height}, width: {shape_width}\n'.
                      format(
                        indent='', 
                        shape_left=shp.left, 
                        shape_height=shp.height, 
                        shape_top=shp.top, 
                        shape_width=shp.width))

                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions'] = OrderedDict()
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['left'] = str(shp.left)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['top'] = str(shp.top)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['width'] = str(shp.width)
                pptx_tree[file_name]['slides'][slide_number]['shapes'][shape_number]['shape dimensions']['height'] = str(shp.height)

    if save_as_json:
        save_json(pptx_tree, file_name+'.json')

    print('Finished')
        
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def read_slide(sld):
    '''
    Takes a slide object and print info regarding the shapes on the given slide. 
    '''

    for x, shp in enumerate(sld.shapes):
        print('{indent:>5}shape index - {x}'.format(indent='', x=x))
        print('{indent:>10}shape name - {shape_name}'.format(indent='', shape_name=shp.name))
        print('{indent:>15}shape type - {shape_type}'.format(indent='', shape_type=shp.shape_type))
        if str(shp.shape_type) == 'PLACEHOLDER (14)':
            print('{indent:>15}placeholder idx - {placeholder_idx}, placeholder type - {placeholder_type}'.
                format(indent='', placeholder_idx=shp.placeholder_format.idx, placeholder_type=shp.placeholder_format.type))
        print('{indent:>15}shape dimensions - left ({shape_left}), top ({shape_top}), height ({shape_height}), width ({shape_width})\n'.
              format(indent='', shape_left=shp.left, shape_top=shp.top, shape_height=shp.height, shape_width=shp.width)) 

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def read_chart_properties(pres_path, slide_num, chart_name):
    '''
    This function prints a given chart's property settings. 

    param: pres_path - full path of target file
    param: slide_num - single slide number 
    param: chart_name - object name as it appears within powerpoint's Object Selection Pane
    '''

    prs = Presentation(pres_path)

    for i, sld in enumerate(prs.slides, start=1):
        if i == slide_num: 
            for x, shp in enumerate(sld.shapes):

                if shp.name == chart_name:

                    print('chart >\n')
                    print('   chart_style: {chart_style}'.format(chart_style=shp.chart.chart_style))
                    print('   has_legend: {has_legend}'.format(has_legend=shp.chart.has_legend))
                    print('   legend: {legend}\n'.format(legend=shp.chart.legend))
                    print('-'*110)
                    
                    caxis = shp.chart.category_axis
                    
                    print('chart > category axis properties\n')
                    print('   has_major_gridlines: {has_major_gridlines}'.format(has_major_gridlines=caxis.has_major_gridlines))
                    print('   has_minor_gridlines: {has_minor_gridlines}'.format(has_minor_gridlines=caxis.has_minor_gridlines))
                    print('   major_tick_mark: {major_tick_mark}'.format(major_tick_mark=caxis.major_tick_mark))
                    print('   maximum_scale: {maximum_scale}'.format(maximum_scale=caxis.maximum_scale))
                    print('   minimum_scale: {minimum_scale}'.format(minimum_scale=caxis.minimum_scale))
                    print('   minor_tick_mark: {minor_tick_mark}'.format(minor_tick_mark=caxis.minor_tick_mark))
                    print('   tick_labels: {tick_labels}'.format(tick_labels=str(caxis.tick_labels)))
                    print('   tick_label_position: {tick_label_position}'.format(tick_label_position=caxis.tick_label_position))
                    print('   tick_labels_font_name: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.name))
                    print('   tick_labels_font_size: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.size))
                    print('   tick_labels_font_bold: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.bold))
                    print('   tick_labels_font_color: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.color))
                    print('   tick_labels_font_italic: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.italic))
                    print('   tick_labels_font_underline: {tick_labels_font}'.format(tick_labels_font=caxis.tick_labels.font.underline))
                    print('   tick_labels_number_format: {tick_labels_number_format}'.format(tick_labels_number_format=caxis.tick_labels.number_format))
                    print('   tick_labels_number_format_is_linked: {tick_labels_number_format_is_linked}'.format(tick_labels_number_format_is_linked=caxis.tick_labels.number_format_is_linked))
                    print('   tick_labels_offset: {tick_labels_offset}'.format(tick_labels_offset=caxis.tick_labels.offset))
                    print('   visible: {visible}\n'.format(visible=caxis.visible))
                    print('-'*110)
                    
                    vaxis = shp.chart.value_axis
                    
                    print('chart > value axis properties\n')
                    print('   has_major_gridlines: {has_major_gridlines}'.format(has_major_gridlines=vaxis.has_major_gridlines))
                    print('   has_minor_gridlines: {has_minor_gridlines}'.format(has_minor_gridlines=vaxis.has_minor_gridlines))
                    print('   major_tick_mark: {major_tick_mark}'.format(major_tick_mark=vaxis.major_tick_mark))
                    print('   maximum_scale: {maximum_scale}'.format(maximum_scale=vaxis.maximum_scale))
                    print('   minimum_scale: {minimum_scale}'.format(minimum_scale=vaxis.minimum_scale))
                    print('   major_unit: {major_unit}'.format(major_unit=vaxis.major_unit))
                    print('   minor_unit: {minor_unit}'.format(minor_unit=vaxis.minor_unit))
                    print('   minor_tick_mark: {minor_tick_mark}'.format(minor_tick_mark=vaxis.minor_tick_mark))
                    print('   tick_labels: {tick_labels}'.format(tick_labels=vaxis.tick_labels))
                    print('   tick_label_position: {tick_label_position}'.format(tick_label_position=vaxis.tick_label_position))
                    print('   tick_labels_font_name: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.name))
                    print('   tick_labels_font_size: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.size))
                    print('   tick_labels_font_bold: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.bold))
                    print('   tick_labels_font_color: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.color))
                    print('   tick_labels_font_italic: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.italic))
                    print('   tick_labels_font_underline: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font.underline))
                    print('   tick_labels_font: {tick_labels_font}'.format(tick_labels_font=vaxis.tick_labels.font))
                    print('   tick_labels_number_format: {tick_labels_number_format}'.format(tick_labels_number_format=vaxis.tick_labels.number_format))
                    print('   tick_labels_number_format_is_linked: {tick_labels_number_format_is_linked}'.format(tick_labels_number_format_is_linked=vaxis.tick_labels.number_format_is_linked))
                    print('   visible: {visible}\n'.format(visible=vaxis.visible))
                    print('-'*110)

                    for item in shp.chart.plots:
                        
                        print('chart > plot\n')
                        print('   plot_categories: {plot_cats}'.format(plot_cats=item.categories))
                        print('   plot_gap_width: {gap_width}'.format(gap_width=item.gap_width))
                        print('   has_data_labels: {has_data_labels}'.format(has_data_labels=item.has_data_labels))
                        print('   overlap: {overlap}'.format(overlap=item.overlap))
                        print('   vary_by_categories: {vary_by_cat}\n'.format(vary_by_cat=item.vary_by_categories))
                        print('-'*110)
                        
                        font = item.data_labels.font
                        
                        print('chart > plot > data labels > font \n')

                        print('   data_label_font_name: {font_name}'.format(font_name=font.name))
                        print('   data_label_font_size: {font_size}'.format(font_size=font.size))
                        print('   data_label_font_bold: {data_label_font}'.format(data_label_font=font.bold))
                        print('   data_label_font_color {font_color}'.format(font_color=font.color))
                        print('   data_label_font_fill {font_fill}'.format(font_fill=font.fill))
                        print('   data_label_font_italic: {font_italic}'.format(font_italic=font.italic))
                        print('   data_label_font_underline: {font_underline}\n'.format(font_underline=font.underline))
                        print('-'*110)
                        
                        for ser in item.series:
                            
                            print('chart > plot > series\n')
                            print('   series_fill_type: {fill_type}'.format(fill_type=ser.fill.type))
                            print('   series_invert_if_neg: {invert_if_neg}'.format(invert_if_neg=ser.invert_if_negative))
                            print('   series_line: {line}'.format(line=ser.line))
                            print('   series_name: {name}'.format(name=ser.name))
                            print('   series_values: {values}'.format(values=ser.values))

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_chart_data_from_prs(pres_path, slide_num, chart_name):
    '''
    This function 1) pulls a given chart's data and 2) returns it as a pandas dataframe object in a list
    
    param: pres_path - full path of target file
    param: slide_num - takes a list of slides 
    param: chart_name - object name as it appears within powerpoint's Object Selection Pane
    '''
    
    prs = Presentation(pres_path)
    
    collection_of_dfs = []

    for i, sld in enumerate(prs.slides, start=1):
        if i in slide_num: 
            for x, shp in enumerate(sld.shapes):
                
                if shp.name == chart_name:
                    
                    plot = shp.chart.plots[0]

                    columns = []
                    data = []
                    for series in plot.series:
                        columns.append(str(series.name))
                        data.append(series.values)
                  
                    data = np.array(data)
                    rows = np.array(plot.categories)
                      
                    df = pd.DataFrame(data.T, index=rows, columns=columns)    
                    collection_of_dfs.append(df)

    return(collection_of_dfs)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def replace_chart_data_in_prs(pres_path, slide_num, chart_name, df):
    '''
    This function 1) enters an existing powerpoint, 2) finds given slides, 3) finds given chart by name and
    4) replaces the given chart's underlying data with new data in the form of a dataframe. 

    param: pres_path - takes the full path of target file
    param: slide_num - takes a list of slides 
    param: chart_name - object name as it appears within powerpoint's Object Selection Pane
    param: df - takes a list of pandas dataframe objects 
    '''
    
    PRES_FOLDER_FOLDER = dirname(pres_path)
    PRES_NAME = basename(pres_path).replace('.pptx','')
    
    prs = Presentation(pres_path)
    
    loop_counter=0

    for i, sld in enumerate(prs.slides, start=1):
        if i in slide_num:
            for x, shp in enumerate(sld.shapes):

                if shp.name == chart_name:

                    single_df = df[loop_counter]
                    chart_data = ChartData()
                    chart_data.categories = single_df.index
 
                    for col_idx, col in enumerate(single_df.columns):
                        chart_data.add_series(col, (single_df.ix[:, col_idx].values))            
                     
                    shp.chart.replace_data(chart_data)
                    
            loop_counter+=1
                
    prs.save('{pres_path}\\{pres_name}_edited.pptx'.format(
            pres_path=PRES_FOLDER_FOLDER, 
            pres_name=PRES_NAME))  

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_slide_layout_names(pptx):
    '''
    Print slide layout names
    '''
    
    for i, slide_layout in enumerate(pptx.slide_layouts):
        print(slide_layout.name)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def return_slide_layout_by_name(pptx, slide_layout_name):
    '''
    Loop over the slide layout object and find slide layout by name, return slide layout 
    object. 

    example: myslide = get_slide_layout_by_name(prs, 'Inhaltsverzeichnis')
             slide = prs.slides.add_slide(myslide)
    '''
    
    for slide_layout in pptx.slide_layouts:
        if slide_layout.name == slide_layout_name:
            return slide_layout
    else:
        raise Exception(
            ('Slide layout: {sld_layout} not found\n').format(
                sld_layout = slide_layout_name))

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_chart_data(shape):

    plot = shape.chart.plots[0]

    columns = []
    data = []
    for series in plot.series:
        columns.append(series.name)
        data.append(series.values)
  
    data = np.array(data)
    rows = np.array(plot.categories)
      
    df = pd.DataFrame(data.T, index=rows, columns=columns)

    return df  

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_chart_data_temp(shape):

    plot = shape.chart.plots[0]

    series_names =[]
    data = []
    for series in plot.series:
        series_names.append(series.name)
        data.append(series.values)

    cols = plot.categories  

    df = pd.DataFrame(data, index=series_names, columns=cols) 

    return df  

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def replace_chart_data(shape, df):

    chart_data = ChartData()
    chart_data.categories = df.index

    for col_idx, col in enumerate(df.columns):
        chart_data.add_series(col, (df.ix[:, col_idx].values))            
     
    shape.chart.replace_data(chart_data)

    return shape

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_slide(pptx, slide_num):
    ''''
    active slides are slides which exist in the VIEW mode,
    not in slide master.
    '''

    return pptx.slides[slide_num]

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_shape(slide_num, shape_name):

    for shp in slide_num.shapes:
        if shp.name == shape_name:
            if shp.is_placeholder:
                p_idx = shp.placeholder_format.idx 
                shp = slide_num.placeholders[p_idx]
            return shp     


'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def copy_txtbox_properties(shp_copy_from, shp_copy_to):
    '''
    Copies over one textbox's properties to another. 
    '''

    # get original slide's shapes dimensions
    left = shp_copy_from.left
    top = shp_copy_from.top
    width = shp_copy_from.width
    height = shp_copy_from.height

    # access textframe property for both original and replica shapes
    txtframe_ori = shp_copy_from.text_frame
    txtframe_rep = shp_copy_to.text_frame 

    # transfer textbox setters from original to replica at textbox level
    txtframe_rep.margin_bottom = txtframe_ori.margin_bottom
    txtframe_rep.margin_left = txtframe_ori.margin_left
    txtframe_rep.margin_right = txtframe_ori.margin_right
    txtframe_rep.margin_top = txtframe_ori.margin_top
    txtframe_rep.vertical_anchor = txtframe_ori.vertical_anchor
    txtframe_rep.word_wrap  = txtframe_ori.word_wrap 
    txtframe_rep.paragraphs[0].text = txtframe_ori.paragraphs[0].text
    txtframe_rep.paragraphs[0].alignment = txtframe_ori.paragraphs[0].alignment

    # color textboxes accordingly
    try:
        color_code = str(shp_copy_from.fill.fore_color.rgb)
        txfill = shp_copy_to.fill
        txfill.solid()
        txfill.fore_color.rgb = RGBColor.from_string(color_code)
    except:
        pass

    # get font size and transfer it to replica shapes
    for paragraph in txtframe_ori.paragraphs:
        for run in paragraph.runs:
            font = run.font
            try:
                font_size = font.size.pt
                t = txtframe_rep.paragraphs[0]
                t.font.size = Pt(font_size)
            except:
                pass

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def copy_chart_properties(shp_copy_from, sld_rep):

    original_shapes_chart_type = str(shp_copy_from.chart.chart_type).split(" ")[0]

    df = get_chart_data(shp_copy_from)
    
    #-------------------------------------------------------------------- 

    add_bar_chart(
        sld=sld_rep, dataframe=df, 
        left=shp_copy_from.left, top=shp_copy_from.top, width=shp_copy_from.width, height=shp_copy_from.height,
        chart_style=shp_copy_from.chart.chart_style, 
        
        has_legend=shp_copy_from.chart.has_legend, 
        legend_position='right',
        legend_in_layout=False,
        legend_horz_offset = 0.1583,
        legend_font_name="Calibri",
        legend_font_size=10,
        legend_font_bold=False,
        legend_font_italic=False,
        legend_font_color=(89,89,89),
        legend_font_brightness=0,
        
        caxis_visible=True,
        caxis_tick_label_position='none',
        caxis_tick_labels_offset=730,
        caxis_has_major_gridlines=shp_copy_from.chart.category_axis.has_major_gridlines,
        caxis_has_minor_gridlines=False,
        caxis_major_tick_mark='outside', 
        caxis_minor_tick_mark='none',
        caxis_tick_labels_font_name="Calibri",
        caxis_tick_labels_font_size=10,
        caxis_tick_labels_font_bold=False,
        caxis_tick_labels_font_italic=False,
        caxis_tick_labels_font_color=(89,89,89),
        
        vaxis_visible=shp_copy_from.chart.value_axis.visible, 
        vaxis_tick_label_position='low',
        vaxis_has_major_gridlines=True,
        vaxis_has_minor_gridlines=False,
        vaxis_major_tick_mark='outside',
        vaxis_minor_tick_mark='none', 
        vaxis_max_scale=100.0, 
        vaxis_min_scale=0,
        vaxis_major_unit=10,
        vaxis_minor_unit=None,
        vaxis_tick_labels_num_format='0"%"', 
        vaxis_tick_labels_font_name="Calibri",
        vaxis_tick_labels_font_bold=True,
        vaxis_tick_labels_font_size=10,
        vaxis_tick_labels_font_italic=False,
        vaxis_tick_labels_font_color=(89,89,89),
        
        plot_has_data_labels=True,
        data_labels_position='outside_end',
        data_labels_num_format='0"%"',
        data_labels_num_format_is_linked=False,
        data_labels_font_name="Calibri",
        data_labels_font_size=9,
        data_labels_font_bold=False,
        data_labels_font_italic=False,
        data_labels_font_color=(0,0,0),
        

        plot_vary_by_cat=False,
        series_color_order='reverse',
        invert_series_color_if_negative=False,
        plot_gap_width=150,
        plot_overlap=-10
        )