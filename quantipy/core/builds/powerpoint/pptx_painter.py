# encoding: utf-8

'''
@author: Majeed.sahebzadha
'''


import copy
import time
import numpy as np
import pandas as pd
import quantipy as qp
from os import path
from collections import OrderedDict
from pptx import Presentation
from quantipy.core.cluster import Cluster
from quantipy.core.chain import Chain
from quantipy.core.helpers.functions import(
    finish_text_key,
    paint_view)
from quantipy.core.builds.powerpoint.add_shapes import(
    chart_selector,
    add_stacked_bar_chart,
    add_textbox,
    add_net)
from quantipy.core.builds.powerpoint.transformations import(
    is_grid_element,
    get_base,
    validate_cluster_orientations,
    drop_hidden_codes,
    partition_view_df,
    strip_html_tags,
    rename_label,
    df_splitter,
    auto_sort,
    round_df_cells)
from quantipy.core.builds.powerpoint.visual_editor import(
    return_slide_layout_by_name)
from pptx.enum.text import(
  PP_ALIGN,
  MSO_AUTO_SIZE,
  MSO_ANCHOR
  )
from pptx.util import(
    Emu,
    Pt,
    Cm,
    Inches
    )
from quantipy.core.builds.powerpoint.add_shapes import (
    percentage_of_num,
    get_cht_plot_height,
    get_upper_cht_plot_gap)

thisdir = path.split(__file__)[0]

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def chain_generator(cluster):
    '''
    Generate chains

    Parameters
    ----------
    cluster : quantipy.Cluster
        quantipy cluster object
    '''

    for chain_name in list(cluster.keys()):
        yield cluster[chain_name]

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def chain_has_weighted_views(chain):
    '''
    Check if a qp.Chain contains weighted frequency views

    Parameters
    ----------
    chain : quantipy.Chain
        quantipy chain object
    '''

    for el in chain.views:
        e0, e1, e2, e3, e4, e5 = el.split('|')
        if e0 == 'x' and e1 == 'f' and e3 == 'y' and e4:
            return True
        elif e2 == 'x++:' and e3 == 'y' and e4:
            return True
    return False

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_grid_el_label(df):
    '''
    Grabs a grid element level label

    Parameters
    ----------
    df : dataframe
        pandas dataframe object
    '''

    grid_el_label = strip_html_tags(df.index[0][0])
    if ' - ' in grid_el_label:
        label = grid_el_label.split(' - ')[-1].strip()
    elif '. ' in grid_el_label:
        label = grid_el_label.split('. ',1)[-1].strip()
    else:
        label = 'Label missing'
    return label

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def df_meta_filter(
    df,
    meta,
    conditions,
    index_key=None):
    '''
    Selects rows based on multiple binary conditions -True/False

    Parameters
    ----------
    df: pandas dataframe
    meta: pandas dataframe
    conditions: dict or pandas.Series object
    index_key: column label or list of column labels / arrays

    example useage: df_meta_filter(
        df,
        meta,
        {'is_pct': True, 'is_weighted': 'True'},
        index_key='label')

    resource: http://stackoverflow.com/questions/34740778/use-series-to-select-rows-from-df-pandas
    http://stackoverflow.com/questions/34726569/get-subsection-of-df-based-on-multiple-conditions
    '''

    con = conditions.copy()

    df = df.reset_index()
    meta = meta.reset_index()

    if not isinstance(con, pd.Series):
        con = pd.Series(con)

    # pull rows where all the conditions are met
    # get subset of df based on labels in conditions
    df = df[(meta == con)[con.index].all(axis=1)]

    if not df.empty:
        if not index_key:
            key_names = ['Values']
        else:
            if not isinstance(index_key, list):
                index_key = [index_key]
            key_names = index_key

        # replace names with labels (note in the future, use text first then labels)
        if len(key_names)>1:
            # use label and overlap those by text which are not empty string
            idx = meta.loc[df.index]['label'].where(meta.loc[df.index]['text']=='', meta.loc[df.index]['text'].values)
        else:
            idx = meta.loc[df.index].set_index(key_names).index

        df = df.set_index(df.columns[0])
        # replace label index with name index
        df.index = idx

        return df
    else:
        # empty df
        return pd.DataFrame()

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def gen_meta_df(painted_df, qp_view):
    '''
    Creates a df containing only metadata

    Parameters
    ----------
    painted_df: pandas dataframe
    unpainted_df: pandas dataframe
    qp_view: quantipy view
    '''

    df_meta = partition_view_df(qp_view.dataframe)[0]
    df_meta['short_name'] = qp_view.meta()['agg']['name']
    df_meta['text'] = qp_view.meta()['agg']['text']
    df_meta['method'] = qp_view.meta()['agg']['method']
    df_meta['is_block'] = qp_view.meta()['agg']['is_block']
    df_meta['is_pct'] = str(qp_view.is_pct())
    df_meta['is_base'] = str(qp_view.is_base())
    df_meta['is_weighted'] = str(qp_view.is_weighted())
    df_meta['is_counts'] = str(qp_view.is_counts())
    df_meta['is_meanstest'] = str(qp_view.is_meanstest())
    df_meta['is_propstest'] = str(qp_view.is_propstest())
    df_meta['is_sum'] = str(qp_view.is_sum())
    df_meta['is_stat'] = str(qp_view.is_stat())
    df_meta['label'] = painted_df.index

    # distinguish between net and expanded
    net_bools=[]
    for row in df_meta.index:
        if qp_view.is_net():
            v_described = qp_view.describe_block()
            all_normal = all(vt == 'normal' for vt in v_described.values())
            if not all_normal:
                if row in v_described:
                    if v_described[row] == 'net':
                        net_bools.append('True')
                    else:
                        net_bools.append('False')
            else:
                net_bools.append('True')
        else:
            net_bools.append('False')
    df_meta['is_net'] = net_bools

    # rearrange the columns
    df_meta = df_meta[['label', 'short_name', 'text', 'method', 'is_pct',
                       'is_net', 'is_weighted', 'is_counts', 'is_block',
                       'is_base', 'is_stat', 'is_sum', 'is_propstest',
                       'is_meanstest']]

    return df_meta

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def same_labels(listofdfs):
    '''
    Before concatenating dfs make sure their row/index labels match.
    Some times, a set of grid element tables which belond to the same
    grid contain different views which is not ideal.

    Parameters
    ----------
    listofdfs: list of pandas dataframes
    '''

    for x in range(0, len(listofdfs)):
        if not all(listofdfs[0].index == listofdfs[x].index):
            raise Exception('index labels mismatch')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def same_num_of_elements(listofdfs):
    '''
    Counts the num of elements in listofdfs, checks if they are all the same.

    Parameters
    ----------
    listofdfs: list of pandas dataframes
    '''

    el_len = [len(el) for el in listofdfs]
    if not all(x == el_len[0] for x in el_len):
        raise Exception('cannot merge elements - uneven '
                        'number of element views.')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def all_same(val_array):
    '''
    Check if all the values in given list the same

    Parameters
    ----------
    numpy_list: numpy array
    '''

    # check if val_array is a numpy array
    if type(val_array).__module__ == np.__name__:
        val = val_array.tolist()
        if isinstance(val[0], list):
            #handle list of lists
            return all(round(x[0]) == round(val[0][0]) for x in val)
        else:
            #handle single list
            return all(round(x) == round(val[0]) for x in val)
    else:
        raise Exception('This function only takes a numpy array')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def insert_values_to_labels(
    add_values_to,
    take_values_from,
    index_position=0):
    '''
    Takes two dfs, adds values from a given row from one df and adds
    it to the other dfs column labels.

    Parameters
    ----------
    add_values_to: pandas dataframe
    take_values_from: pandas dataframe
    index_position: int, optional
    '''

    # check 1 - if the labels from both dfs are the same
    if all(add_values_to.columns == take_values_from.columns):

        # pull a given row's values
        row_vals = take_values_from.ix[[index_position],:].values
        # flatten the list of values
        row_vals = row_vals.flatten()
        # get column labels
        col_labels = add_values_to.columns

        # loop over and add the values to the labels
        for x,y in zip(col_labels, row_vals):
            col_name = x + " (n=" + str(int(round(y))) +")"
            add_values_to.rename(columns={x: col_name}, inplace=True)

        return add_values_to
    else:
        raise Exception('Cannot add values to df labels')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def PowerPointPainter(
    path_pptx,
    meta,
    cluster,
    path_pptx_template=None,
    slide_layout='Blank',
    text_key=None,
    force_chart=True,
    force_crossbreak=None,
    base_type='weighted',
    base_repr=None,
    include_nets=True,
    shape_properties=None,
    display_var_names=True,
    date_range=None,
    split_busy_dfs=False,
    verbose=True):
    '''
    Builds PowerPoint file (PPTX) from cluster, list of clusters, or
    dictionary of clusters.

    Parameters
    ----------
    path_pptx : str
        PowerPoint file path
    meta : dict
        metadata as dictionary used to paint datframes
    cluster : quantipy.Cluster / list / dict
        container for cluster(s)
    path_pptx_template : str, optional
        full path to PowerPoint template
    slide_layout : str / int, optional
        valid slide layout name or index
    text_key : str, optional
        language
    force_chart : boolean, optional
        ues default settings to produce a PowerPoint file
    force_crossbreak : str / list, optional
        use given crossbreaks to build a PowerPoint file
    base_type : str, optional
        use weighted or unweighted base
    include_nets : str / boolean
        True/False: include/exclude net views in chart data
        'partly': include nets in chart data except for array summary charts
    shape_properties : dict, optional
        keys as format properties, values as change from default
    display_var_names : boolean
        variable names append to question labels
    split_busy_dfs : boolean
        if True, spreads busy dataframes evenly across multiple slide
    '''
    if verbose:
        print((
            '\n{ast}\n{ast}\n{ast}\nINITIALIZING POWERPOINT '
            'AUTOMATION SCRIPT...'.format(ast='*' * 80)))

    # check path extension
    if path_pptx.endswith('.pptx'):
        path_pptx = path_pptx[:-5]
    elif path_pptx.endswith('/') or path_pptx.endswith('\\'):
        raise Exception('File name not provided')

    # check base type string
    base_type = base_type.lower()

    # render cluster
    names = []
    clusters = []
    if isinstance(cluster, Cluster):
        names.append(cluster.name)
        clusters.append(cluster)
    elif isinstance(cluster, list):
        for c in cluster:
            names.append(c.name)
            clusters.append(c)
    elif isinstance(cluster, dict):
        names_clusters_dict = cluster
        for sheet_name, c in cluster.items():
            names.append(sheet_name)
            clusters.append(c)

    # default settings
    default_props = {
        'crossbreak': ['@'],
        'chart_type': 'bar',
        'sort_order': 'none',
        'chart_color': 'green',
        'fixed_categories': [],
        'base_description': '',
        'chart_layout': '1',
        'slide_title_text': 'Click to add title',
        'question_label': 'Unknown',
        'copied_from': '',
        'center_header': '',
        'right_footer': '',
        'title_footer': ''}

    spec = meta['sets'].get('spec', False)

    # update 'crossbreak' key's value in default_props if
    # force_crossbreak parameter is true
    if force_crossbreak:
        if isinstance(force_crossbreak, list):
            pass
        elif isinstance(force_crossbreak, str):
            force_crossbreak = [force_crossbreak]
        for c in force_crossbreak:
            default_props['crossbreak'].append(c)

    if not path_pptx_template:
        path_pptx_template = path.join(
            thisdir,
            'templates\default_template.pptx')

    # get the default text key if none provided
    if text_key is None:
        text_key = finish_text_key(meta, text_key)

    # default shape properties (minimum level, only shape dimensions)
    # if none provided
    if shape_properties is None:
         shape_properties = {
            'header_shape': {
                'left': 284400,
                'top': 1007999,
                'width': 8582400,
                'height': 468000},
            'chart_shape': {
                'bar': {
                    'left': 284400,
                    'top': 1475999,
                    'width': 8582400,
                    'height': 4140000},
                'stacked_bar': {
                    'left': 284400,
                    'top': 1475999,
                    'width': 8582400,
                    'height': 4140000},
                'column': {
                    'left': 284400,
                    'top': 1475999,
                    'width': 8582400,
                    'height': 4140000},
                'pie': {
                    'left': 284400,
                    'top': 1475999,
                    'width': 8582400,
                    'height': 4140000},
                'line': {
                    'left': 284400,
                    'top': 1475999,
                    'width': 8582400,
                    'height': 4140000}},
            'footer_shape': {
                'left': 284400,
                'top': 5652000,
                'width': 8582400,
                'height': 396000}}

    ############################################################################
    ############################################################################
    ############################################################################

    # loop over clusters, returns pptx for each cluster
    for cluster_name, cluster in zip(names, clusters):
        if verbose:
            print((
                '\nPowerPoint minions are building your PPTX, '
                'please stand by...\n\n{indent:>2}Building '
                'PPTX for {file_name}').format(
                    indent='',
                    file_name=cluster_name))

        # log start time
        pptx_start_time = time.time()

        # check if cluster is empty
        if not cluster:
            raise Exception("'{}' cluster is empty".format(cluster_name))

        # ensure all chains in cluster have the same orientation
        validate_cluster_orientations(cluster)

        # pull orientation of chains in cluster
        orientation = cluster[list(cluster.keys())[0]].orientation

        # open pptx template file
        prs = Presentation(path_pptx_template)

        # log slide number
        slide_num = len(prs.slides)

        # Get Client and Topic tag from meta    
        if isinstance(spec, dict):
            topic = "{}".format(spec.get('topic', ""))        
            client = "{}".format(spec.get('name', ""))            
        else:
            topic = ""        
            client = ""                    
        
        ############################################################################
        # frontpage title ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
        ############################################################################

        title_shape=shape_properties['title_shape']
        client_date_shape=shape_properties['client_date_shape']

        if title_shape['addtext_frontpage']:

            for shape in prs.slides[0].shapes:
                if shape.name == title_shape['shape_name_frontpage']:
                    shape.text = topic

        if client_date_shape['addtext']:

            for shape in prs.slides[0].shapes:
                if shape.name == title_shape['shape_name']:
                    shape.text = client_date_shape['t_d_text'].format(client,date_range)

        ############################################################################
        # X ORIENTATION CODE ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
        ############################################################################

        if orientation == 'x':

            # grid element storage dict
            grid_container = []
            # translated views contains names of all views
            # which have been translated
            translated_views = []

            # This section tries to finds, pull and build grid element
            # dataframes by matching the downbreak name against the grid element name.
            # Each downbreak is therefore checked against all keys in masks.
            for chain in chain_generator(cluster):

                # list of crossbreak name
                crossbreaks = chain.content_of_axis
                # single downbreak name
                downbreak = chain.source_name

                '----CHART AND BASE DATA CONDITIONS ----------------------------'

                # table selection conditions for chart shape
                chartdata_conditions = OrderedDict([
                    ('is_pct', 'True'),
                    ('is_weighted', 'True'),
                    ('is_sum', 'False')])

                # Net settings
                net_setup = shape_properties.get('net_setup', False)
                if not include_nets:
                    chartdata_conditions.update({'is_net': 'False'})
                    chartdata_conditions_grid = copy.deepcopy(chartdata_conditions)
                elif include_nets == True:
                    if net_setup:
                        chartdata_conditions.update({'is_net': 'False'})
                        chartdata_conditions_grid = copy.deepcopy(chartdata_conditions)
                    else:
                        chartdata_conditions_grid = copy.deepcopy(chartdata_conditions)
                #elif include_net == 'partly':
                else:
                    chartdata_conditions_grid = copy.deepcopy(chartdata_conditions)
                    chartdata_conditions_grid.update({'is_net': 'False'})


                #---------------------------------------------------------------
                # table selection conditions for footer/base shape
                base_conditions = OrderedDict([
                    ('is_base', 'True'),
                    ('short_name', 'cbase'),
                    ('is_weighted', 'True' if base_type == 'weighted' else 'False')])

                '----PULL METADATA DETAILS -------------------------------------'

                # for each downbreak, try and pull it's meta
                if force_chart:
                    meta_props = []
                else:
                    if downbreak in meta['columns']:
                        if 'properties' in meta['columns'][downbreak]:
                            meta_props = meta['columns'][downbreak]['properties']
                        else:
                            meta_props = []
                    else:
                        meta_props = []

                if text_key['x'] in list(meta['columns'][downbreak]['text'].keys()):
                    question_label = meta['columns'][downbreak]['text'][text_key['x']]
                else:
                    question_label = list(meta['columns'][downbreak]['text'].values())[0]
                chart_type = meta_props['chart_type'] if 'chart_type' in meta_props else default_props['chart_type']
                layout_type = meta_props['chart_layout'] if 'chart_layout' in meta_props else default_props['chart_layout']
                sort_order = meta_props['sort_order'] if 'sort_order' in meta_props else default_props['sort_order']
                fixed_categories = meta_props['fixed_categories'] if 'fixed_categories' in meta_props else default_props['fixed_categories']
                if fixed_categories:
                    fixed_values = [int(x['value']) for x in fixed_categories]
                    values = loc_values = meta['columns'][downbreak]['values']
                    if isinstance(loc_values, str):
                        loc_values = loc_values.split('@')
                        values = meta[loc_values.pop(0)]
                        while loc_values:
                            values = values[loc_values.pop(0)]
                    fixed_categories = [
                        item['text'][text_key['x']]
                        for item in values
                        if item['value'] in fixed_values
                    ]
                slide_title_text = meta_props['slide_title'] if 'slide_title' in meta_props else default_props['slide_title_text']
                copied_from = meta_props['copied_from'] if 'copied_from' in meta_props else default_props['copied_from']
                base_description = meta_props['base_text'] if 'base_text' in meta_props else default_props['base_description']

                '----IF GRID THEN-----------------------------------------------'

                # loop over items in masks
                for grid in meta['masks']:
                    for x in range(0, len(meta['masks'][grid]['items'])):
                        gridname = meta['masks'][grid]['items'][x]['source'].split('columns@')[-1]
                        if downbreak == gridname:

                            if text_key['x'] in list(meta['masks'][grid]['text'].keys()):
                                grid_question_label = meta['masks'][grid]['text'][text_key['x']]
                            else:
                                grid_question_label = list(meta['masks'][grid]['text'].values())[0]

                            # check if grid is in grid container, if it's not then continue
                            if not grid in grid_container:
                                grid_container += [grid]
                                remaining_elements = [
                                    grid_element['source'].split('@')[1]
                                    for grid_element in meta['masks'][grid]['items'][0:]]

                                '----GROUP GRID-CHAIN VIEWS-------------------------------------'

                                grouped_grid_views = []

                                for grid_element_name in remaining_elements:
                                    grid_chain = cluster[grid_element_name]

                                    #prepare grid element labels
                                    grid_el_label = meta['columns'][grid_element_name]['text'][text_key['x']]

                                    if grid_el_label.startswith(grid_question_label):
                                        grid_el_label = grid_el_label.split(grid_question_label)[-1].strip()
                                        if grid_el_label.startswith('- '):
                                            grid_el_label = grid_el_label[2:]

                                    # use weighted freq views if available
                                    has_weighted_views = chain_has_weighted_views(grid_chain)

                                    #if the conditions for base and chartdata's "is_weighted" key
                                    #is True but there are no weighted views in the chain then use
                                    #unweighted views
                                    if not has_weighted_views:
                                        if chartdata_conditions['is_weighted']=='True':
                                            chartdata_conditions['is_weighted'] = 'False'
                                            chartdata_conditions_grid['is_weighted'] = 'False'
                                            #an unweighted chart can only have unweighted base
                                            if base_conditions['is_weighted']=='True':
                                                base_conditions['is_weighted'] = 'False'

                                    views_on_chain = []
                                    meta_on_g_chain = []

                                    # loop over views in chain
                                    for v in grid_chain.views:
                                        dk = grid_chain.data_key
                                        fk = grid_chain.filter

                                        # only pull '@' based views as these will be concatenated together
                                        view = grid_chain[dk][fk][grid_element_name]['@'][v]

                                        view.translate_metric(
                                            text_key['x'][0],
                                            set_value='meta')

                                        trans_var_name = '{}x@'.format(grid_chain.name)
                                        if not trans_var_name in translated_views:
                                            translated_views.append(trans_var_name)

                                        # paint view
                                        df = paint_view(meta, view)
                                        # flatten df
                                        df = partition_view_df(df)[0]
                                        # get meta data
                                        df_meta = gen_meta_df(df, view)
                                        # append
                                        meta_on_g_chain.append(df_meta)
                                        views_on_chain.append(df)

                                    # this var will be overwritten but its okay for now.
                                    grped_g_meta = pd.concat(
                                        meta_on_g_chain,
                                        axis=0)

                                    # concat all the views together on a single chain
                                    mdf = pd.concat(views_on_chain, axis=0)
                                    mdf.rename(
                                        columns={mdf.columns[0]: grid_el_label},
                                        inplace=True)
                                    grouped_grid_views.append(mdf)

                                '----CONCAT AND PREPARE GRID-CHAIN VIEWS------------------------'

                                # before merging all grid elements together, 2 checks are carried out:
                                # 1. ensure all grid elements have the same number of views
                                same_num_of_elements(grouped_grid_views)
                                # 2. ensure all the grid index labels are the name
                                same_labels(grouped_grid_views)

                                # concat all grid chains in grouped_grid_views together
                                merged_grid_df = pd.concat(grouped_grid_views, axis=1)
                                merged_grid_df = merged_grid_df.fillna(0.0)

                                slide_num += 1
                                if verbose:
                                    print((
                                        '\n{indent:>5}Slide {num}. '
                                        'Adding a 100% STACKED BAR CHART '
                                        'for {qname} cut by '
                                        'Total{war_msg}'.format(
                                            indent='',
                                            num=slide_num,
                                            qname=grid,
                                            war_msg='')))

                                #extract df for net
                                if net_setup:
                                    net_setup_stacked_bar = net_setup.get('stacked_bar', False)
                                    if net_setup_stacked_bar:
                                        df_grid_table_net = df_meta_filter(
                                            merged_grid_df,
                                            grped_g_meta,
                                            {'is_net' : 'True'},
                                            index_key='label')
                                            
                                #extract df for chart
                                df_grid_table = df_meta_filter(
                                    merged_grid_df,
                                    grped_g_meta,
                                    chartdata_conditions_grid,
                                    index_key='label')

                                #extract df for base
                                df_grid_base = df_meta_filter(
                                    merged_grid_df,
                                    grped_g_meta,
                                    base_conditions,
                                    index_key='text')

                                if not df_grid_table.empty:

                                    # if not all the values in the grid's df are the same
                                    # then add the values to the grids column labels
                                    if not all_same(df_grid_base.values[0]):
                                        df_grid_table = insert_values_to_labels(
                                            df_grid_table,
                                            df_grid_base,
                                            index_position=0)
                                        if base_description:
                                            #remove the word "Base:" from the description
                                            description = base_description.split(': ')[-1]
                                            #grab the label for base from the df
                                            base_label = df_grid_base.index[0]
                                            #put them together
                                            base_text = '{}: {}'.format(
                                                base_label,
                                                description)
                                        else:
                                            base_text = ''
                                    else:
                                        base_text = get_base(
                                            df_grid_base,
                                            base_description,
                                            True)

                                    if base_repr and ('Base' in base_text):
                                        base_text = base_text.replace('Base', base_repr)

                                    # get question label
                                    if display_var_names:
                                        if shape_properties['short_grid_name']:
                                            grid_label = grid.partition('.')[0]
                                        else:
                                            grid_label = grid
                                        grid_question_label = '{}. {}'.format(
                                            grid_label,
                                            strip_html_tags(grid_question_label))

                                    # format table values
                                    df_grid_table = np.round(df_grid_table/100, 4)

                                    '----ADDPEND SLIDE TO PRES--------------------------------------'

                                    if isinstance(slide_layout, int):
                                        slide_layout_obj = prs.slide_layouts[slide_layout]
                                    else:
                                        slide_layout_obj = return_slide_layout_by_name(
                                            prs,
                                            slide_layout)

                                    slide = prs.slides.add_slide(slide_layout_obj)

                                    '----ADD SHAPES TO SLIDE----------------------------------------'

                                    ''' title shape'''
                                    if title_shape['addtext']:

                                        for shape in slide.placeholders:
                                            if shape.name == title_shape['shape_name']:
                                                shape.text = topic

                                    ''' header shape '''
                                    sub_title_shp = add_textbox(
                                        slide,
                                        text=grid_question_label,
                                        **(shape_properties['header_shape']
                                            if shape_properties else {}))

                                    ''' net table '''
                                    if include_nets and net_setup:
                                        save_width = shape_properties['chart_shape']['stacked_bar']['width']
                                        if net_setup_stacked_bar['show_table']:
                                            if not df_grid_table_net.empty:
                                                df_grid_table_net = round_df_cells(df_grid_table_net,
                                                                                   net_setup_stacked_bar['table_decimals'])
                                                if net_setup_stacked_bar['add_percent_sign']:
                                                    df_grid_table_net = df_grid_table_net.astype(str) + '%'
                                                cols = len(df_grid_table_net.T.columns)
                                                shapes=shape_properties['chart_shape']['stacked_bar']
                                                shapes['legend_position']='bottom'
                                                shapes['width'] -= net_setup_stacked_bar['table_column_width'] * cols
                                                # Set net table size and position
                                                height = shapes['height']
                                                top = shapes['top']
                                                left = shapes['left'] + shapes['width']
                                                width = net_setup_stacked_bar['table_column_width']
                                                net_table = add_net(slide, df_grid_table_net.T, height=height, width=width, top=top, left=left)                                            
                                            
                                    ''' chart shape '''
                                    chart_shp = chart_selector(
                                        slide,
                                        df_grid_table,
                                        chart_type='stacked_bar',
                                        **(shape_properties['chart_shape']['stacked_bar']
                                            if shape_properties else {}))

                                    if include_nets and net_setup:
                                        shape_properties['chart_shape']['stacked_bar']['width'] = save_width
                                            
                                    ''' footer shape '''
                                    if base_text:
                                        base_text_shp = add_textbox(
                                            slide,
                                            text=base_text,
                                            **(shape_properties['footer_shape']
                                                if shape_properties else {}))

                '----IF NOT GRID THEN-------------------------------------------'

                if 'crossbreak' in meta_props:
                    if meta_props['crossbreak'] != '@':
                        target_crossbreaks = default_props['crossbreak'] + meta_props['crossbreak'].split(',')
                    else:
                        target_crossbreaks = meta_props['crossbreak'].split(',')
                else:
                    target_crossbreaks = default_props['crossbreak']

                for crossbreak in crossbreaks:
                    if crossbreak in target_crossbreaks:

                        '----GROUP NON GRID-CHAIN VIEWS---------------------------------'

                        # are there any weighted views in this chain?
                        has_weighted_views = chain_has_weighted_views(chain)

                        # if "is_weighted" is True but there's no weighted views
                        # use unweighted views
                        if not has_weighted_views:
                            if chartdata_conditions['is_weighted']=='True':
                                chartdata_conditions['is_weighted'] = 'False'
                                # an unweighted chart can only have unweighted base
                                if base_conditions['is_weighted'] == 'True':
                                    base_conditions['is_weighted'] = 'False'

                        views_on_chain = []
                        meta_on_chain = []

                        for v in chain.views:
                            dk = chain.data_key
                            fk = chain.filter

                            view = chain[dk][fk][downbreak][crossbreak][v]

                            trans_var_name = '{}x{}'.format(
                                downbreak,
                                crossbreak)

                            if trans_var_name not in translated_views:
                                view.translate_metric(
                                    text_key['x'][0],
                                    set_value='meta')

                            # paint view
                            df = paint_view(meta, view)
                            # flatten df
                            df = partition_view_df(df)[0]
                            # get meta data
                            df_meta = gen_meta_df(df, view)
                            # append to vars
                            meta_on_chain.append(df_meta)
                            views_on_chain.append(df)

                        '----CONCAT AND PREPARE NON GRID-CHAIN VIEWS--------------------'

                        grped_meta = pd.concat(meta_on_chain, axis=0)
                        grped_df = pd.concat(views_on_chain, axis=0)
                        grped_df = grped_df.fillna(0.0)

                        # replace '@' with 'Total'
                        grped_df = rename_label(
                            grped_df,
                            '@',
                            'Total',
                            orientation='Top')

                        # extract df for net
                        if net_setup:
                            df_table_net = df_meta_filter(
                            grped_df,
                            grped_meta,
                            {'is_net': 'True'},
                            index_key='label')
                            # standardise table values
                            df_table_net = np.round(df_table_net.fillna(0.0) / 100, 4)
                            
                        #extract df for chart
                        df_table = df_meta_filter(
                            grped_df,
                            grped_meta,
                            chartdata_conditions,
                            index_key='label')

                        #extract df for base
                        df_base = df_meta_filter(
                            grped_df,
                            grped_meta,
                            base_conditions,
                            index_key='text')

                        if not df_table.empty:

                            # append nets to fixed categories
                            for x, item in enumerate(grped_meta['is_net'].tolist()):
                                if item == 'True':
                                    if fixed_categories!=[]:
                                        fixed_categories.append(grped_meta['label'].tolist()[x])
                                    else:
                                        fixed_categories = [grped_meta['label'].tolist()[x]]

                            # sort df whilst excluding fixed cats
                            if sort_order == 'ascending':
                                df_table = auto_sort(
                                    df=df_table,
                                    fixed_categories=fixed_categories,
                                    column_position=0,
                                    ascend=True)
                            elif sort_order == 'descending':
                                df_table = auto_sort(
                                    df=df_table,
                                    fixed_categories=fixed_categories,
                                    column_position=0,
                                    ascend=False)

                            # if not all the values in the grid's df are the same
                            # then add the values to the grids column labels
                            if not all_same(df_base.values):
                                df_table = insert_values_to_labels(
                                    df_table,
                                    df_base,
                                    index_position=0)
                                base_text = base_description
                            else:
                                if not df_base.empty:
                                    base_text = get_base(
                                        df_base,
                                        base_description,
                                        False)
                                else:
                                    raise Exception('Base dataframe empty for "{}".'.format(downbreak))

                            if base_repr and ('Base' in base_text):
                                base_text = base_text.replace('Base', base_repr)

                            # standardise table values
                            df_table = np.round(df_table.fillna(0.0)/100, 4)

                            # get question label
                            if display_var_names:
                                if shape_properties['short_grid_name'] and '_grid' in downbreak:
                                    downbreak_label = downbreak.partition('{')[2].partition('}')[0]
                                else:
                                    downbreak_label = downbreak
                                question_label = '{}. {}'.format(
                                    downbreak_label,
                                    strip_html_tags(question_label))


                            # handle incorrect chart type assignment
                            if len(df_table.index) > 15 and chart_type == 'pie':
                                chart_type='bar'

                            '----SPLIT DFS & LOOP OVER THEM---------------------------------'

                            if split_busy_dfs:
                                # split large dataframes
                                collection_of_dfs = df_splitter(
                                    df_table,
                                    min_rows=5,
                                    max_rows=15)
                            else:
                                # dont split large/busy dataframes
                                collection_of_dfs = [df_table]

                            for i, df_table_slice in enumerate(collection_of_dfs):

                                '----ADDPEND SLIDE TO PRES--------------------------------------'

                                if isinstance(slide_layout, int):
                                    slide_layout_obj = prs.slide_layouts[slide_layout]
                                else:
                                    slide_layout_obj = return_slide_layout_by_name(
                                        prs,
                                        slide_layout)

                                slide = prs.slides.add_slide(slide_layout_obj)

                                '----ADD SHAPES TO SLIDE----------------------------------------'

                                ''' title shape'''
                                if title_shape['addtext']:

                                    for shape in slide.placeholders:
                                        if shape.name == title_shape['shape_name']:
                                            shape.text = topic

                                ''' title shape '''
                                if i > 0:
                                    cont_question_label = '{} (continued {})'.format(
                                        question_label,
                                        i+1)
                                else:
                                    cont_question_label = question_label

                                ''' header shape '''
                                sub_title_shp = add_textbox(
                                    slide,
                                    text=cont_question_label,
                                    **(shape_properties['header_shape']
                                        if shape_properties else {}))

                                ''' chart shape '''
                                numofcols = len(df_table_slice.columns)
                                numofrows = len(df_table_slice.index)

                                # handle incorrect chart type assignment
                                if chart_type == 'pie' and numofcols > 1:
                                    chart_type = 'bar'

                                # turn legend off if table contains 1 series unless its a pie chart
                                if numofcols == 1:
                                    legend_switch = False
                                    if chart_type == 'pie':
                                        legend_switch = True
                                else:
                                    legend_switch = True

                                if 'has_legend' in shape_properties['chart_shape'][chart_type]:
                                    shape_properties['chart_shape'][chart_type]['has_legend'] = legend_switch

                                # Net settings
                                if include_nets and net_setup:
                                    net_setup = net_setup.get(chart_type, False)
                                    if not net_setup == False and net_setup['show_nets']:
                                        if len(collection_of_dfs) == 1:
                                            if not df_table_net.empty:
                                                if net_setup['separator']:
                                                    df_table_slice = df_table_slice.T
                                                    df_table_slice.insert(len(df_table_slice.columns), 'net_separator', 1.01)
                                                    df_table_slice = df_table_slice.T
                                                #df_table_slice.loc[len(df_table_slice)]=0
                                                df_table_slice = pd.concat([df_table_slice, df_table_net])
                                                shape_properties['chart_shape']['bar']['separator_color'] = net_setup['separator_color']                                    
                                    
                                chart = chart_selector(
                                    slide,
                                    df_table_slice,
                                    chart_type=chart_type,
                                     **(shape_properties['chart_shape'][chart_type]
                                        if shape_properties else {}))
                                
                                ''' footer shape '''
                                base_text_shp = add_textbox(
                                    slide,
                                    text=base_text,
                                    **(shape_properties['footer_shape']
                                        if shape_properties else {}))

                                slide_num += 1
                                if verbose:
                                    print((
                                        '\n{indent:>5}Slide {slide_number}. '
                                        'Adding a {chart_name} '
                                        'CHART for {question_name} '
                                        'cut by {crossbreak_name} '
                                        '{x}'.format(
                                            indent='',
                                            slide_number=slide_num,
                                            chart_name=chart_type.upper().strip(),
                                            question_name=downbreak,
                                            crossbreak_name='Total' if crossbreak == '@' else crossbreak,
                                            x='(cont ('+str(i)+'))' if i > 0 else '')))

                        else:
                            if verbose:
                                print((
                                    '\n{indent:>5}***Skipping {question_name}, '
                                    'no views match your conditions: '
                                    '{conditions}'.format(
                                        indent='',
                                        question_name=downbreak,
                                        conditions=chartdata_conditions)))

            prs.save('{}.pptx'.format(path_pptx))
            print('Created: {}.pptx\n'.format(path_pptx))

        ############################################################################
        # Y ORIENTATION CODE ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
        ############################################################################

        if orientation == 'y':

            # raise error is cluster is y orientated
            raise TypeError('y orientation not supported yet')

    #-------------------------------------------------------------------------
    #-------------------------------------------------------------------------
    #-------------------------------------------------------------------------

    if verbose:
        pptx_elapsed_time = time.time() - pptx_start_time
        print((
            '\n{indent:>2}Presentation saved, '
            'time elapsed: {time:.2f} seconds\n'
            '\n{line}'.format(
                indent='',
                time=pptx_elapsed_time,
                line= '_' * 80)))
