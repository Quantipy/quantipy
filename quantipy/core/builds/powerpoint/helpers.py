
from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import numpy as np
import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def ChartData_from_DataFrame(df, number_format="0%", xl_number_format='0.00%'):
    """
    Return a CategoryChartData instance from the given Pandas DataFrame.
    """

    def get_parent(sub_categories, line, pos):
        """
        Return the sub_category's parent given its lineage position.
        """

        for subcat in sub_categories:
            if subcat.label == line[pos]:
                return subcat


    cd = CategoryChartData(number_format=number_format)

    if isinstance(df.index, pd.MultiIndex):
        cats = []
        for line in df.index.unique().tolist():
            for l, lvl in enumerate(line):
                if l == 0:
                    if not any([lvl == cat.label for cat in cats]):
                        cats.append(cd.add_category(lvl))
                else:
                    parent = get_parent(cats, line, 0)
                    if l > 1:
                        for i in range(1, l):
                            parent = get_parent(parent.sub_categories, line, i)
                    sub_categories = parent.sub_categories
                    seen = [lvl == subcat.label for subcat in sub_categories]
                    if not any(seen):
                        parent.add_sub_category(lvl)
    else:
        categories = tuple(df.index.values.tolist())
        cd.categories = categories

    for col in df.columns:
        values = [
            value if value==value else None
            for value in df[col].values.tolist()
        ]
        if isinstance(col,tuple):
            series = (" - ".join(col), tuple(values))
        else:
            series = (col, tuple(values))

        cd.add_series(*series, number_format=xl_number_format)

    return cd


def DataFrame_from_ChartData(cd):
    """
    Return a Pandas DataFrame from the given CategoryChartData instance.
    """

    tuples = get_category_tuples(cd.categories)
    if isinstance(tuples[0], tuple):
        idx = pd.MultiIndex.from_tuples(tuples)
    else:
        idx = pd.Index(tuples)
    cols = pd.Index([series.name for series in cd])
    data = list(zip(*[
        [value if value != None else np.NaN for value in series.values]
        for series in cd
    ]))

    df = pd.DataFrame(data, index=idx, columns=cols)

    return df


def get_category_tuples(categories):
    """
    Return the hierarchical tuples representation of the given categories.
    """

    def get_sub_category_tuples(sub_categories, records, record):
        """
        Update records with sub_caregories' hierarhical tuples.
        """

        if sub_categories:
            for sub_category in sub_categories:
                sub_record = record[:] + [sub_category.label]
                sub_sub_categories = sub_category.sub_categories
                get_sub_category_tuples(sub_sub_categories, records, sub_record)
        else:
            records.append(tuple(record))

    records = []
    for category in categories:
        record = [category.label]
        get_sub_category_tuples(category.sub_categories, records, record)

    if all([len(record) == 1 for record in records]):
        records = [record[0] for record in records]

    return tuple(records)


def DataFrame_from_Chart(chart):
    """
    Return a Pandas DataFrame from an existing Chart.
    """

    plot = chart.plots[0]
    idx = [series.name for series in plot.series]
    if plot.categories.depth > 1:
        cols = pd.MultiIndex.from_tuples(plot.categories.flattened_labels)
    else:
        cols = pd.Index(plot.categories)
    data = [series.values for series in plot.series]
    df = pd.DataFrame(data, index=idx, columns=cols).T

    return df


def verify_ChartData_vs_DataFrame(cd, df):
    """
    Print a comparison of the given ChartData and DataFrame.
    """

    print((df.fillna('') == DataFrame_from_ChartData(cd).fillna('')))


def verify_DataFrame_vs_DataFrame(df1, df2):
    """
    Print a comparison of the two given DataFrames.
    """

    print((df1.fillna('') == df2.fillna('')))

def example_dataframe(hierarchical):
    """
    Return an example Pandas DataFrame.
    """

    if hierarchical:
        idx = pd.MultiIndex.from_tuples((
                ('USA', 'CA', 'San Francisco'),
                ('USA', 'CA', 'Los Angeles'),
                ('USA', 'NY', 'New York'),
                ('USA', 'NY', 'Albany'),
                ('CAN', 'AL', 'Calgary'),
                ('CAN', 'AL', 'Edmunton'),
                ('CAN', 'ON', 'Toronto'),
                ('CAN', 'ON', 'Ottawa'),
            )
        )
    else:
        idx = pd.Index((
                'San Francisco',
                'Los Angeles',
                'New York',
                'Albany',
                'Calgary',
                'Edmunton',
                'Toronto',
                'Ottawa',
            )
        )

    cols = pd.Index(['Series 1', 'Series 2', 'Series 3', 'Series 4'])
    data = [
        (4, 8, 1, 5),
        (4, 7, 2, 3),
        (4, 6, np.NaN, 2),
        (4, 5, 4, np.NaN),
        (4, 4, 5, 7),
        (4, np.NaN, 6, 1),
        (4, 2, 7, 4),
        (np.NaN, 1, 8, 3)
    ]

    df = pd.DataFrame.from_records(data, index=idx, columns=cols)

    return df


def add_slide_with_chart(prs, cd, title):
    """
    Adds a new slide with a chart on it using the given chart data.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
    x, y, cx, cy = 1524000, 1397000, 6096000, 4064000
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        x, y, cx, cy,
        cd
    ).chart

    txBox = slide.shapes.add_textbox(x, y-1000000, cx, cy)
    txBox.text_frame.text = title

    return chart
